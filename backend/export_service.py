"""
Export Service - Generate downloadable exports (PPTX, PDF, Word)
"""
import os
import io
import json
from typing import Optional, List
from pydantic import BaseModel

# ============================================================================
# Models
# ============================================================================

class Slide(BaseModel):
    title: str
    content: Optional[str] = None
    bullet_points: Optional[List[str]] = None
    speaker_notes: Optional[str] = None
    layout: Optional[str] = "title_content"
    image_url: Optional[str] = None

class ExportSlidesRequest(BaseModel):
    slides: List[Slide]
    title: str
    format: str = "pptx"  # pptx, pdf, html
    template: Optional[str] = None

class ExportWordRequest(BaseModel):
    title: str
    sections: List[dict]
    include_toc: bool = True

# ============================================================================
# PowerPoint Export
# ============================================================================

async def export_slides_pptx(request: ExportSlidesRequest) -> bytes:
    """Generate a PowerPoint presentation optimized for designer editing

    Creates a clean, well-structured presentation with:
    - Clear content hierarchy
    - Comprehensive speaker notes
    - Visual suggestions in notes
    - Professional formatting ready for design
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RgbColor
        from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
    except ImportError:
        raise ValueError("python-pptx not installed. Run: pip install python-pptx")

    prs = Presentation()
    prs.slide_width = Inches(10)  # Standard 16:9 ratio
    prs.slide_height = Inches(7.5)

    # Title slide - clean and professional
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = request.title

    # Style title for designer-friendly output
    if title.text_frame.paragraphs:
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    # Content slides - optimized for clarity and designer editing
    for slide_data in request.slides:
        # Use blank layout for maximum designer flexibility
        blank_layout = prs.slide_layouts[6]  # Blank slide
        slide = prs.slides.add_slide(blank_layout)

        # Add title text box
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(1)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.word_wrap = True

        title_p = title_frame.paragraphs[0]
        title_p.text = slide_data.title
        title_p.font.size = Pt(32)
        title_p.font.bold = True
        title_p.font.color.rgb = RgbColor(0, 0, 0)

        # Add subtitle if present
        content_top = Inches(1.8)
        if slide_data.content and not slide_data.bullet_points:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.5))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.word_wrap = True

            subtitle_p = subtitle_frame.paragraphs[0]
            subtitle_p.text = slide_data.content[:200]  # First 200 chars as subtitle
            subtitle_p.font.size = Pt(18)
            subtitle_p.font.color.rgb = RgbColor(100, 100, 100)
            content_top = Inches(2.2)

        # Add bullet points with proper formatting
        if slide_data.bullet_points:
            content_box = slide.shapes.add_textbox(Inches(1), content_top, Inches(8), Inches(4))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            for i, point in enumerate(slide_data.bullet_points):
                if i == 0:
                    p = content_frame.paragraphs[0]
                else:
                    p = content_frame.add_paragraph()

                # Clean bullet point text
                clean_point = point.strip()
                if clean_point.startswith(('•', '-', '*')):
                    clean_point = clean_point[1:].strip()

                p.text = clean_point
                p.level = 0
                p.font.size = Pt(20)
                p.font.color.rgb = RgbColor(50, 50, 50)
                p.space_before = Pt(12)

                # Add bullet
                p.font.name = 'Arial'

        # Add comprehensive speaker notes for designers
        if slide_data.speaker_notes or slide_data.layout or slide_data.image_url:
            notes_slide = slide.notes_slide
            notes_text = []

            # Add speaker notes
            if slide_data.speaker_notes:
                notes_text.append("CONTENT NOTES:")
                notes_text.append(slide_data.speaker_notes)
                notes_text.append("")

            # Add layout suggestion
            if slide_data.layout:
                notes_text.append(f"SUGGESTED LAYOUT: {slide_data.layout}")
                notes_text.append("")

            # Add image suggestion
            if slide_data.image_url or (slide_data.bullet_points and len(slide_data.bullet_points) > 3):
                notes_text.append("DESIGN SUGGESTIONS:")
                if slide_data.image_url:
                    notes_text.append(f"- Suggested image: {slide_data.image_url}")
                notes_text.append("- Consider adding visual elements (icons, illustrations, or photos)")
                notes_text.append("- Use whitespace effectively - avoid text-heavy slides")
                notes_text.append("- Apply brand colors and fonts")
                notes_text.append("")

            # Add content structure note
            if slide_data.bullet_points:
                notes_text.append(f"CONTENT STRUCTURE: {len(slide_data.bullet_points)} key points")
                notes_text.append("- Each point should be visually distinct")
                notes_text.append("- Consider using icons or numbers for each point")

            notes_slide.notes_text_frame.text = "\n".join(notes_text)

    # Add a final "Designer Notes" slide
    notes_layout = prs.slide_layouts[6]
    final_slide = prs.slides.add_slide(notes_layout)

    notes_box = final_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
    notes_frame = notes_box.text_frame
    notes_frame.word_wrap = True

    designer_notes = [
        "DESIGNER NOTES",
        "",
        "This presentation has been structured for professional design enhancement.",
        "",
        "Content Structure:",
        f"  • {len(request.slides)} content slides",
        "  • Clear hierarchy: titles, subtitles, bullet points",
        "  • Speaker notes contain detailed guidance",
        "",
        "Design Recommendations:",
        "  • Apply consistent brand colors and fonts",
        "  • Add relevant images, icons, or illustrations",
        "  • Ensure proper whitespace and visual balance",
        "  • Use animations sparingly and purposefully",
        "  • Maintain readability (minimum 18pt body text)",
        "",
        "All content is ready for your creative enhancement!",
    ]

    for i, line in enumerate(designer_notes):
        if i == 0:
            p = notes_frame.paragraphs[0]
            p.text = line
            p.font.size = Pt(28)
            p.font.bold = True
        else:
            p = notes_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(14) if line.startswith("  •") else Pt(16)
            if line and not line.startswith(" "):
                p.font.bold = True

    # Save to bytes
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()

# ============================================================================
# Word Export
# ============================================================================

async def export_word(request: ExportWordRequest) -> bytes:
    """Generate a Word document"""
    try:
        from docx import Document
        from docx.shared import Inches, Pt
    except ImportError:
        raise ValueError("python-docx not installed. Run: pip install python-docx")
    
    doc = Document()
    
    # Title
    doc.add_heading(request.title, 0)
    
    # Table of contents placeholder
    if request.include_toc:
        doc.add_paragraph("Table of Contents", style='Heading 1')
        doc.add_paragraph("[Table of contents will be generated by Word]")
        doc.add_page_break()
    
    # Sections
    for section in request.sections:
        doc.add_heading(section.get('title', 'Section'), level=1)
        
        if section.get('content'):
            doc.add_paragraph(section['content'])
        
        if section.get('subsections'):
            for sub in section['subsections']:
                doc.add_heading(sub.get('title', 'Subsection'), level=2)
                if sub.get('content'):
                    doc.add_paragraph(sub['content'])
    
    # Save to bytes
    output = io.BytesIO()
    doc.save(output)
    output.seek(0)
    return output.read()

# ============================================================================
# HTML Export (for PDF conversion)
# ============================================================================

async def export_slides_html(request: ExportSlidesRequest) -> str:
    """Generate HTML representation of slides"""
    html_parts = [f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>{request.title}</title>
    <style>
        body {{ font-family: 'Segoe UI', Arial, sans-serif; margin: 0; padding: 20px; }}
        .slide {{ 
            page-break-after: always; 
            border: 1px solid #ddd; 
            padding: 40px; 
            margin-bottom: 20px;
            min-height: 500px;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            color: white;
            border-radius: 8px;
        }}
        .slide h1 {{ font-size: 2.5em; margin-bottom: 30px; }}
        .slide ul {{ font-size: 1.3em; line-height: 1.8; }}
        .slide p {{ font-size: 1.2em; line-height: 1.6; }}
        .notes {{ 
            background: #f5f5f5; 
            padding: 15px; 
            margin-top: 20px; 
            border-left: 4px solid #667eea;
            color: #333;
            border-radius: 4px;
        }}
    </style>
</head>
<body>
    <div class="slide">
        <h1>{request.title}</h1>
    </div>
"""]
    
    for slide in request.slides:
        html_parts.append(f'<div class="slide">')
        html_parts.append(f'<h1>{slide.title}</h1>')
        
        if slide.bullet_points:
            html_parts.append('<ul>')
            for point in slide.bullet_points:
                html_parts.append(f'<li>{point}</li>')
            html_parts.append('</ul>')
        elif slide.content:
            html_parts.append(f'<p>{slide.content}</p>')
        
        if slide.speaker_notes:
            html_parts.append(f'<div class="notes"><strong>Speaker Notes:</strong> {slide.speaker_notes}</div>')
        
        html_parts.append('</div>')
    
    html_parts.append('</body></html>')
    return '\n'.join(html_parts)
